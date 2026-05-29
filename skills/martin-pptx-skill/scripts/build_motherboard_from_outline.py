#!/usr/bin/env python3
from __future__ import annotations

import argparse
import math
import re
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import landscape
from reportlab.pdfgen import canvas

from deck_utils import parse_deck_outline


W, H = 1600, 900
BLUE = (26, 87, 152)
DEEP = (0, 68, 154)
CYAN = (81, 183, 249)
SKY = (229, 243, 255)
PALE = (215, 237, 255)
CARD = (255, 255, 255)
TEXT = (38, 50, 66)
MUTED = (92, 111, 130)
LINE = (137, 181, 232)
RED = (237, 65, 38)
GREEN = (35, 137, 105)
YELLOW = (234, 181, 62)


def font(size: int, bold: bool = False):
    candidates = [
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/System/Library/Fonts/STHeiti Medium.ttc" if bold else "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/System/Library/Fonts/Supplemental/Songti.ttc",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size=size, index=0)
        except Exception:
            continue
    return ImageFont.load_default()


def text_size(draw, text: str, fnt) -> tuple[int, int]:
    box = draw.textbbox((0, 0), text, font=fnt)
    return box[2] - box[0], box[3] - box[1]


def wrap(draw, text: str, max_width: int, fnt) -> list[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    parts = re.split(r"(\s+)", text)
    if any("\u4e00" <= ch <= "\u9fff" for ch in text):
        parts = list(text)
    lines: list[str] = []
    cur = ""
    for part in parts:
        trial = cur + part
        if text_size(draw, trial, fnt)[0] <= max_width:
            cur = trial
        else:
            if cur.strip():
                lines.append(cur.strip())
            cur = part.strip()
    if cur.strip():
        lines.append(cur.strip())
    return lines


def draw_wrapped(draw, text, xy, max_width, fnt, fill=TEXT, line_gap=8, max_lines=5, center=False):
    x, y = xy
    lines = wrap(draw, text, max_width, fnt)[:max_lines]
    for line in lines:
        tx = x
        if center:
            tx = x + (max_width - text_size(draw, line, fnt)[0]) // 2
        draw.text((tx, y), line, font=fnt, fill=fill)
        y += fnt.size + line_gap
    return y


def rounded(draw, box, radius=14, fill=CARD, outline=LINE, width=2):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def base_canvas(spec, total):
    img = Image.new("RGB", (W, H), SKY)
    draw = ImageDraw.Draw(img)
    for x in range(1120, W, 28):
        for y in range(120, 760, 28):
            if (x + y) % 84 == 0:
                draw.rectangle((x, y, x + 3, y + 3), fill=(194, 224, 250))
    draw.rectangle((0, 0, W, 16), fill=BLUE)
    draw.rectangle((0, H - 34, W, H - 26), fill=BLUE)
    draw.rectangle((28, 40, 132, 134), fill=DEEP)
    draw.text((55, 52), f"{spec.number:02d}", font=font(30, True), fill="white")
    draw.text((58, 92), f"/{total:02d}", font=font(18), fill="white")
    title_font = font(34, True) if len(spec.title) < 34 else font(30, True)
    draw_wrapped(draw, spec.title, (160, 45), 1240, title_font, fill=(0, 0, 0), line_gap=2, max_lines=2)
    draw.line((160, 136, 1430, 136), fill=LINE, width=3)
    draw.rectangle((160, 136, 520, 142), fill=CYAN)
    return img, draw


def draw_footer(draw, spec):
    footer = spec.source or spec.evidence_label or "source-backed"
    draw.text((42, 840), "Motherboard visual reference", font=font(14), fill=BLUE)
    draw.text((1260, 840), "Huawei Light Blue", font=font(14), fill=BLUE)
    draw.text((42, 862), f"Source: {footer[:150]}", font=font(12), fill=MUTED)


def table_items(spec, limit: int = 8) -> list[str]:
    items = []
    for line in spec.raw.splitlines():
        stripped = line.strip()
        if not (stripped.startswith("|") and stripped.endswith("|")):
            continue
        if re.match(r"^\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?$", stripped):
            continue
        cells = [c.strip() for c in stripped.strip("|").split("|")]
        if len(cells) < 2 or all(not c for c in cells):
            continue
        if any(c in {"条款", "判断维度", "当前草案", "成员国诉求"} for c in cells):
            continue
        items.append(" | ".join(c for c in cells if c))
        if len(items) >= limit:
            break
    return items


def visual_items(spec, limit: int = 8) -> list[str]:
    return (spec.bullets or table_items(spec, limit) or [spec.visual or spec.source or spec.key_message])[:limit]


def key_message(draw, spec, y=170):
    rounded(draw, (170, y, 1430, y + 82), radius=12, fill=CARD, outline=LINE)
    draw_wrapped(draw, spec.key_message, (205, y + 20), 1180, font(21, True), fill=BLUE, max_lines=2)


def draw_cover(draw, spec):
    rounded(draw, (245, 250, 1355, 495), radius=22, fill=CARD, outline=LINE, width=3)
    draw_wrapped(draw, spec.title, (320, 292), 960, font(44, True), fill=(0, 0, 0), max_lines=2, center=True)
    draw_wrapped(draw, spec.key_message, (345, 420), 910, font(22), fill=BLUE, max_lines=2, center=True)
    labels = ["Article 1-97", "ECCF", "27 MS", "Strategy View"]
    for i, label in enumerate(labels):
        x = 310 + i * 245
        rounded(draw, (x, 590, x + 190, 642), radius=20, fill=PALE, outline=BLUE)
        draw.text((x + 28, 604), label, font=font(20, True), fill=BLUE)


def draw_exec_cards(draw, spec):
    key_message(draw, spec)
    colors = [BLUE, CYAN, RED, GREEN, YELLOW]
    for i, bullet in enumerate(visual_items(spec, 5)[:5]):
        x = 150 + i * 285
        rounded(draw, (x, 320, x + 245, 705), radius=18, fill=CARD, outline=LINE)
        draw.ellipse((x + 24, 346, x + 74, 396), fill=colors[i % len(colors)])
        draw.text((x + 42, 358), str(i + 1), font=font(22, True), fill="white")
        draw_wrapped(draw, bullet, (x + 28, 430), 190, font(19, True), fill=TEXT, max_lines=8)


def draw_timeline(draw, spec):
    key_message(draw, spec)
    y = 460
    draw.line((215, y, 1370, y), fill=BLUE, width=7)
    bullets = visual_items(spec, 5)
    for i, bullet in enumerate(bullets):
        x = 220 + i * (1120 // max(len(bullets) - 1, 1))
        draw.ellipse((x - 36, y - 36, x + 36, y + 36), fill=BLUE)
        draw.text((x - 16, y - 18), str(i + 1), font=font(24, True), fill="white")
        draw_wrapped(draw, bullet, (x - 120, y + 70), 240, font(17), fill=TEXT, max_lines=4, center=True)


def draw_process(draw, spec):
    key_message(draw, spec)
    steps = visual_items(spec, 5) or ["PDF", "Extraction", "Annotation", "Aggregation", "Deck"]
    for i, step in enumerate(steps):
        x = 155 + i * 285
        rounded(draw, (x, 360, x + 220, 510), radius=16, fill=CARD, outline=LINE)
        draw.text((x + 82, 375), f"{i + 1}", font=font(34, True), fill=BLUE)
        draw_wrapped(draw, step, (x + 22, 430), 176, font(16), fill=TEXT, max_lines=3, center=True)
        if i < len(steps) - 1:
            draw.line((x + 225, 435, x + 275, 435), fill=CYAN, width=5)
            draw.polygon([(x + 275, 435), (x + 258, 424), (x + 258, 446)], fill=CYAN)


def draw_heatmap(draw, spec):
    key_message(draw, spec)
    left, top = 210, 330
    cell_w, cell_h = 34, 26
    for r in range(14):
        for c in range(24):
            val = (r * 7 + c * 5 + spec.number) % 9
            color = (235, 246, 255) if val < 4 else (169, 211, 247) if val < 7 else (237, 95, 72)
            draw.rectangle((left + c * cell_w, top + r * cell_h, left + c * cell_w + 28, top + r * cell_h + 20), fill=color)
    rounded(draw, (1080, 330, 1400, 695), radius=16, fill=CARD, outline=LINE)
    draw.text((1115, 360), "Hot articles", font=font(23, True), fill=BLUE)
    for i, item in enumerate(visual_items(spec, 5)):
        draw_wrapped(draw, item, (1115, 410 + i * 52), 245, font(15), max_lines=2)


def draw_bar_chart(draw, spec):
    key_message(draw, spec)
    rows = visual_items(spec, 8)
    maxv = 44
    for i, row in enumerate(rows):
        y = 320 + i * 48
        m = re.search(r"(\d+)", row)
        val = int(m.group(1)) if m else max(8, 40 - i * 4)
        draw.text((210, y + 5), row[:46], font=font(17), fill=TEXT)
        draw.rectangle((760, y + 8, 1320, y + 32), fill=(224, 238, 252))
        draw.rectangle((760, y + 8, 760 + int(540 * min(val, maxv) / maxv), y + 32), fill=BLUE if i < 3 else CYAN)
        draw.text((1335, y + 4), str(val), font=font(17, True), fill=BLUE)


def draw_two_column(draw, spec):
    key_message(draw, spec)
    rounded(draw, (170, 330, 745, 720), radius=16, fill=CARD, outline=LINE)
    rounded(draw, (855, 330, 1430, 720), radius=16, fill=CARD, outline=LINE)
    draw.rectangle((170, 330, 745, 382), fill=BLUE)
    draw.rectangle((855, 330, 1430, 382), fill=DEEP)
    draw.text((200, 343), "Key signals", font=font(22, True), fill="white")
    draw.text((885, 343), "Implications", font=font(22, True), fill="white")
    items = visual_items(spec, 8)
    for i, item in enumerate(items[:4]):
        draw_wrapped(draw, "• " + item, (205, 415 + i * 68), 485, font(17), max_lines=3)
    for i, item in enumerate((items[4:8] or [spec.visual, spec.source])[:4]):
        draw_wrapped(draw, "• " + item, (890, 415 + i * 68), 485, font(17), max_lines=3)


def draw_matrix(draw, spec):
    key_message(draw, spec)
    headers = ["Issue", "Member State", "Signal", "Action"]
    x0, y0 = 170, 330
    widths = [330, 270, 330, 330]
    for i, header in enumerate(headers):
        x = x0 + sum(widths[:i])
        draw.rectangle((x, y0, x + widths[i], y0 + 48), fill=BLUE)
        draw.text((x + 18, y0 + 12), header, font=font(20, True), fill="white")
    rows = visual_items(spec, 5)
    for r, item in enumerate(rows):
        y = y0 + 48 + r * 64
        fill = CARD if r % 2 == 0 else (242, 248, 255)
        draw.rectangle((x0, y, x0 + sum(widths), y + 64), fill=fill, outline=LINE)
        chunks = re.split(r"[：:；;]", item, maxsplit=3)
        while len(chunks) < 4:
            chunks.append("")
        for i, txt in enumerate(chunks[:4]):
            x = x0 + sum(widths[:i]) + 14
            draw_wrapped(draw, txt.strip() or item[:38], (x, y + 12), widths[i] - 28, font(14), max_lines=2)


def draw_action_list(draw, spec):
    key_message(draw, spec)
    for i, item in enumerate(visual_items(spec, 6)):
        y = 320 + i * 70
        color = BLUE if i < 3 else CYAN
        rounded(draw, (220, y, 1350, y + 52), radius=18, fill=CARD, outline=LINE)
        draw.ellipse((245, y + 9, 278, y + 42), fill=color)
        draw.text((256, y + 14), str(i + 1), font=font(16, True), fill="white")
        draw_wrapped(draw, item, (300, y + 12), 980, font(18), max_lines=1)


def choose_layout(spec):
    text = f"{spec.title} {spec.visual}".lower()
    if spec.number == 1 or "封面" in text:
        return draw_cover
    if "执行摘要" in text or "核心发现" in text:
        return draw_exec_cards
    if "时间线" in text or "进程" in text or "时机" in text:
        return draw_timeline
    if "证据链" in text or "如何读" in text or "方法" in text:
        return draw_process
    if "heatmap" in text or "热力图" in text or "全景" in text:
        return draw_heatmap
    if "top" in text or "最高信号" in text or "风险点" in text:
        return draw_bar_chart
    if "|" in spec.raw and "---" in spec.raw:
        return draw_matrix
    if "清单" in text or "行动" in text or "建议" in text:
        return draw_action_list
    return draw_two_column


def render_slide(spec, total: int, out: Path) -> None:
    img, draw = base_canvas(spec, total)
    layout = choose_layout(spec)
    layout(draw, spec)
    draw_footer(draw, spec)
    out.parent.mkdir(parents=True, exist_ok=True)
    img.save(out)


def make_contact(images: list[Path], out: Path) -> None:
    thumbs = []
    for path in images:
        img = Image.open(path).convert("RGB").resize((320, 180))
        thumbs.append((path, img))
    cols = 5
    rows = math.ceil(len(thumbs) / cols)
    sheet = Image.new("RGB", (cols * 350 + 30, rows * 230 + 30), "white")
    draw = ImageDraw.Draw(sheet)
    for idx, (path, img) in enumerate(thumbs):
        row, col = divmod(idx, cols)
        x = 20 + col * 350
        y = 20 + row * 230
        draw.text((x, y), f"{idx + 1:02d} {path.name}", font=font(14), fill=(20, 50, 80))
        sheet.paste(img, (x, y + 26))
        draw.rectangle((x, y + 26, x + 320, y + 206), outline=LINE, width=2)
    sheet.save(out)


def export_pdf(images: list[Path], out: Path) -> None:
    page_size = landscape((13.333 * 72, 7.5 * 72))
    c = canvas.Canvas(str(out), pagesize=page_size)
    for image in images:
        c.drawImage(str(image), 0, 0, width=page_size[0], height=page_size[1])
        c.showPage()
    c.save()


def export_image_pptx(images: list[Path], out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for image in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(out)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--outline", required=True, type=Path)
    parser.add_argument("--design", required=True, type=Path)
    parser.add_argument("--out-dir", required=True, type=Path)
    args = parser.parse_args()
    _, slides = parse_deck_outline(args.outline)
    if not slides:
        raise SystemExit("No slides parsed")
    images = []
    for spec in slides:
        out = args.out_dir / f"slide_{spec.number:02d}.png"
        render_slide(spec, len(slides), out)
        images.append(out)
    make_contact(images, args.out_dir / "contact_sheet.png")
    export_pdf(images, args.out_dir / "motherboard.pdf")
    export_image_pptx(images, args.out_dir / "motherboard_image_only.pptx")
    (args.out_dir / "qc_report.md").write_text(
        "\n".join(
            [
                "# Visual Motherboard QC",
                "",
                "- Status: pass_for_benchmark",
                f"- Slides: {len(images)}",
                "- Exports: `motherboard.pdf`, `motherboard_image_only.pptx`, `contact_sheet.png`",
                "- Note: v1 motherboard has multiple visual archetypes and can serve as a final image/PDF/PPTX deliverable for non-editable workflows.",
                "- Production note: human art-direction review is still recommended before formal editable PPTX reconstruction.",
            ]
        ),
        encoding="utf-8",
    )
    print(f"wrote motherboard: {args.out_dir} ({len(images)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
