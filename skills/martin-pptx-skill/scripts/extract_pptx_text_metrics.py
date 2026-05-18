#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return (shape.text or "").strip()


def iter_font_sizes(shape) -> list[float]:
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(round(run.font.size.pt, 2))
    return sizes


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx", required=True, type=Path)
    parser.add_argument("--out", required=True, type=Path)
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    slides = []
    all_sizes: list[float] = []
    total_text_shapes = 0
    total_image_shapes = 0

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        slide_sizes: list[float] = []
        image_shapes = 0
        for shape in slide.shapes:
            text = shape_text(shape)
            if text:
                total_text_shapes += 1
                sizes = iter_font_sizes(shape)
                slide_sizes.extend(sizes)
                texts.append(
                    {
                        "text": text,
                        "font_sizes": sizes,
                        "left": round(shape.left.inches, 3),
                        "top": round(shape.top.inches, 3),
                        "width": round(shape.width.inches, 3),
                        "height": round(shape.height.inches, 3),
                    }
                )
            if shape.shape_type == 13:
                image_shapes += 1
        total_image_shapes += image_shapes
        all_sizes.extend(slide_sizes)
        notes = ""
        if getattr(slide, "has_notes_slide", False):
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes = ""
        slides.append(
            {
                "slide": idx,
                "texts": texts,
                "visible_text": "\n".join(t["text"] for t in texts),
                "text_shape_count": len(texts),
                "image_shape_count": image_shapes,
                "notes": notes,
                "min_font_size": min(slide_sizes) if slide_sizes else None,
                "max_font_size": max(slide_sizes) if slide_sizes else None,
            }
        )

    result = {
        "pptx": str(args.pptx),
        "summary": {
            "slide_count": len(prs.slides),
            "total_text_shapes": total_text_shapes,
            "total_image_shapes": total_image_shapes,
            "min_font_size": min(all_sizes) if all_sizes else None,
            "max_font_size": max(all_sizes) if all_sizes else None,
        },
        "slides": slides,
    }
    args.out.parent.mkdir(parents=True, exist_ok=True)
    args.out.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(result["summary"], ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
