#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from deck_utils import parse_deck_outline


POLICY_BLUE = RGBColor(0x1A, 0x57, 0x98)
SKY = RGBColor(0xE6, 0xF3, 0xFF)
PALE = RGBColor(0xD7, 0xED, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
BODY = RGBColor(0x38, 0x46, 0x53)
RED = RGBColor(0xED, 0x41, 0x26)
GREEN = RGBColor(0x2C, 0x8A, 0x6A)
LINE = RGBColor(0x88, 0xB5, 0xE8)


def add_box(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(0.8)
    return box


def add_text(slide, text, x, y, w, h, size=16, color=BODY, bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei", fill=None, line=None):
    if fill is None:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        shape = add_box(slide, x, y, w, h, fill, line or fill, radius=True)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.07)
    tf.margin_right = Inches(0.07)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_bullets(slide, bullets, x, y, w, h, size=15):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size if len(bullet) < 90 else max(12, size - 2))
        p.font.color.rgb = BODY
        p.space_after = Pt(4)
    return shape


def add_background(slide, background: Path | None):
    if background and background.exists():
        slide.shapes.add_picture(str(background), 0, 0, width=Inches(13.333), height=Inches(7.5))
        return
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = SKY
    add_box(slide, 0.0, 0.0, 13.333, 0.18, POLICY_BLUE, POLICY_BLUE)
    add_box(slide, 0.0, 7.22, 13.333, 0.06, POLICY_BLUE, POLICY_BLUE)


def add_header(slide, idx, total, title, subtitle=""):
    add_text(slide, f"{idx:02d}\n/{total:02d}", 0.15, 0.32, 0.92, 0.76, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font="Arial", fill=POLICY_BLUE, line=POLICY_BLUE)
    title_size = 28 if len(title) <= 34 else 26 if len(title) <= 46 else 24
    add_text(slide, title, 1.32, 0.24, 11.35, 0.72, size=title_size, color=BLACK, bold=True)
    if subtitle:
        sub_size = 15 if len(subtitle) < 80 else 13
        add_text(slide, subtitle, 1.32, 1.02, 11.05, 0.42, size=sub_size, color=BLACK)


def add_footer(slide, spec):
    label = spec.evidence_label or "Source-backed"
    add_text(slide, f"Evidence: {label}", 0.42, 6.90, 4.5, 0.30, size=10, color=POLICY_BLUE, font="Calibri")
    if spec.source:
        add_text(slide, f"Source: {spec.source[:130]}", 5.1, 6.90, 7.35, 0.30, size=10, color=BODY, font="Calibri")


def choose_layout(spec):
    raw = spec.raw
    if "|" in raw and "---" in raw:
        return "table"
    if len(spec.bullets) >= 6:
        return "dense_bullets"
    if len(spec.bullets) <= 3:
        return "message_cards"
    return "standard"


def add_standard_content(slide, spec, layout):
    add_text(slide, spec.key_message, 0.86, 1.72, 11.65, 0.78, size=18, color=POLICY_BLUE, bold=True, fill=WHITE, line=LINE)
    if layout == "message_cards":
        y = 3.0
        items = spec.bullets or [spec.visual or spec.source or "待补充可视化要点"]
        for i, item in enumerate(items[:3]):
            x = 0.92 + i * 4.18
            add_text(slide, str(i + 1), x, y, 0.42, 0.42, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font="Arial", fill=POLICY_BLUE, line=POLICY_BLUE)
            add_text(slide, item, x + 0.55, y - 0.08, 3.35, 1.28, size=15, color=BODY, fill=WHITE, line=LINE)
    elif layout == "table":
        add_text(slide, "条款 / 议题 / 策略价值", 0.92, 2.88, 11.5, 0.38, size=16, color=WHITE, bold=True, fill=POLICY_BLUE, line=POLICY_BLUE)
        add_bullets(slide, spec.bullets[:6], 1.05, 3.4, 11.1, 2.65, size=13)
    elif layout == "dense_bullets":
        add_box(slide, 0.86, 2.82, 11.7, 3.35, WHITE, LINE, radius=True)
        add_bullets(slide, spec.bullets[:8], 1.05, 2.98, 11.2, 2.95, size=13)
    else:
        left = spec.bullets[:3]
        right = spec.bullets[3:6]
        add_text(slide, "Key signals", 0.9, 2.9, 5.6, 0.38, size=16, color=WHITE, bold=True, fill=POLICY_BLUE, line=POLICY_BLUE)
        add_bullets(slide, left, 1.05, 3.38, 5.25, 2.45, size=14)
        add_text(slide, "Implications", 6.85, 2.9, 5.55, 0.38, size=16, color=WHITE, bold=True, fill=POLICY_BLUE, line=POLICY_BLUE)
        add_bullets(slide, right or [spec.visual or "保留为 speaker note / visual reconstruction hint"], 7.0, 3.38, 5.2, 2.45, size=14)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--outline", required=True, type=Path)
    parser.add_argument("--design", required=True, type=Path)
    parser.add_argument("--motherboard-dir", type=Path)
    parser.add_argument("--background", type=Path)
    parser.add_argument("--out", required=True, type=Path)
    parser.add_argument("--run-id", required=True)
    parser.add_argument("--max-slides", type=int, default=0)
    args = parser.parse_args()

    meta, slides = parse_deck_outline(args.outline)
    if args.motherboard_dir:
        contact_sheet = args.motherboard_dir / "contact_sheet.png"
        slide_refs = sorted(args.motherboard_dir.glob("slide_*.png"))
        if not contact_sheet.exists() or not slide_refs:
            raise SystemExit(f"Motherboard missing or incomplete: {args.motherboard_dir}")
    if args.max_slides:
        slides = slides[: args.max_slides]
    if not slides:
        raise SystemExit("No slides parsed from outline")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = len(slides)
    for idx, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        add_background(slide, args.background)
        subtitle = spec.key_message if idx == 1 else ""
        add_header(slide, idx, total, spec.title, subtitle if idx == 1 else "")
        if idx == 1:
            add_text(slide, meta.get("subtitle", "Benchmark deck"), 1.3, 2.18, 10.7, 0.48, size=20, color=POLICY_BLUE, bold=True, align=PP_ALIGN.CENTER, fill=WHITE, line=LINE)
            add_text(slide, spec.key_message, 1.58, 3.05, 10.2, 1.28, size=18, color=BODY, align=PP_ALIGN.CENTER, fill=WHITE, line=LINE)
            add_bullets(slide, spec.bullets[:4], 2.1, 4.82, 9.35, 1.12, size=13)
        else:
            add_standard_content(slide, spec, choose_layout(spec))
        add_footer(slide, spec)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            notes.text = f"Run: {args.run_id}\nSource: {spec.source}\nVisual hint: {spec.visual}\nKey message: {spec.key_message}"

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(f"wrote {args.out} ({total} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
